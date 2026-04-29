import json
import uuid
from datetime import datetime
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, BackgroundTasks
from fastapi.responses import StreamingResponse
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession
from io import BytesIO

from app.database import get_db, AsyncSessionLocal
from app.auth.entra import get_current_user
from app.auth.rbac import require_role, Role
from app.models.slides import PPTTemplate
from app.models.resource import UserRole
from app.schemas.slides import PPTTemplateOut, GenerateSlidesRequest, JobStatus
from app.services.ppt_service import extract_placeholders, populate_template
from app.services.blob_service import upload_blob, download_blob

router = APIRouter()

# In-memory job store for MVP
_jobs: dict[str, dict] = {}


@router.get("/templates", response_model=list[PPTTemplateOut])
async def list_templates(db: AsyncSession = Depends(get_db), user: UserRole = Depends(get_current_user)):
    res = await db.execute(select(PPTTemplate))
    return res.scalars().all()


@router.post("/templates", response_model=PPTTemplateOut)
async def upload_template(
    file: UploadFile = File(...),
    name: str = "",
    tags: str = "",
    db: AsyncSession = Depends(get_db),
    user: UserRole = Depends(require_role(Role.PM)),
):
    contents = await file.read()
    placeholders = extract_placeholders(contents)
    template_id = uuid.uuid4()
    blob_path = f"templates/{template_id}.pptx"
    await upload_blob(blob_path, contents)

    from pptx import Presentation
    slide_count = len(Presentation(BytesIO(contents)).slides)

    tpl = PPTTemplate(
        id=template_id,
        name=name or file.filename or "Template",
        tags=tags,
        slide_count=slide_count,
        blob_path=blob_path,
        placeholder_map=json.dumps(placeholders),
        uploaded_by=user.entra_oid,
    )
    db.add(tpl)
    await db.commit()
    await db.refresh(tpl)
    return tpl


@router.get("/templates/{template_id}", response_model=PPTTemplateOut)
async def get_template(template_id: uuid.UUID, db: AsyncSession = Depends(get_db),
                       user: UserRole = Depends(get_current_user)):
    t = await db.get(PPTTemplate, template_id)
    if not t:
        raise HTTPException(404, "Template not found")
    return t


@router.delete("/templates/{template_id}")
async def delete_template(
    template_id: uuid.UUID, db: AsyncSession = Depends(get_db),
    user: UserRole = Depends(require_role(Role.PM)),
):
    t = await db.get(PPTTemplate, template_id)
    if not t:
        raise HTTPException(404, "Template not found")
    await db.delete(t)
    await db.commit()
    return {"deleted": True}


async def _run_generation_job(job_id: str, request: GenerateSlidesRequest):
    _jobs[job_id]["status"] = "running"
    try:
        async with AsyncSessionLocal() as db:
            tpl = await db.get(PPTTemplate, request.template_id)
            if not tpl:
                raise ValueError("Template not found")
            template_bytes = await download_blob(tpl.blob_path)
            placeholders = json.loads(tpl.placeholder_map)

            from app.agents.slide_agent import resolve_placeholders
            data_map = await resolve_placeholders(placeholders, request, db)

            output = populate_template(template_bytes, data_map)
            out_path = f"generated/{job_id}.pptx"
            await upload_blob(out_path, output)

            tpl.last_used_at = datetime.utcnow()
            await db.commit()

        _jobs[job_id]["status"] = "completed"
        _jobs[job_id]["blob_path"] = out_path
    except Exception as e:
        _jobs[job_id]["status"] = "failed"
        _jobs[job_id]["error"] = str(e)


@router.post("/generate", response_model=JobStatus)
async def generate_slides(
    data: GenerateSlidesRequest,
    background: BackgroundTasks,
    db: AsyncSession = Depends(get_db),
    user: UserRole = Depends(get_current_user),
):
    job_id = str(uuid.uuid4())
    _jobs[job_id] = {"status": "queued", "user": user.entra_oid}
    background.add_task(_run_generation_job, job_id, data)
    return JobStatus(job_id=job_id, status="queued")


@router.get("/jobs/{job_id}", response_model=JobStatus)
async def get_job(job_id: str, user: UserRole = Depends(get_current_user)):
    job = _jobs.get(job_id)
    if not job:
        raise HTTPException(404, "Job not found")
    return JobStatus(
        job_id=job_id,
        status=job["status"],
        download_url=f"/api/slides/jobs/{job_id}/download" if job["status"] == "completed" else None,
        error=job.get("error"),
    )


@router.get("/jobs/{job_id}/download")
async def download_job(job_id: str, user: UserRole = Depends(get_current_user)):
    job = _jobs.get(job_id)
    if not job or job["status"] != "completed":
        raise HTTPException(404, "Job not ready")
    data = await download_blob(job["blob_path"])
    return StreamingResponse(
        BytesIO(data),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="generated_{job_id}.pptx"'},
    )
