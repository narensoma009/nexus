import re
from io import BytesIO
from pptx import Presentation

TOKEN_PATTERN = re.compile(r"\{\{([A-Z_]+)\}\}")


def infer_token_type(name: str) -> str:
    if any(k in name for k in ["TABLE", "ROSTER", "LIST", "SPREAD", "MILESTONE"]):
        return "table"
    if any(k in name for k in ["CHART", "TREND", "HEATMAP"]):
        return "chart"
    if any(k in name for k in ["SUMMARY", "NARRATIVE", "NOTES", "INSIGHTS",
                               "RISKS", "DECISIONS", "ACTIONS"]):
        return "ai"
    if any(k in name for k in ["DATE", "PERIOD", "QUARTER"]):
        return "auto"
    return "data"


def extract_placeholders(pptx_bytes: bytes) -> list[dict]:
    prs = Presentation(BytesIO(pptx_bytes))
    found: dict[str, str] = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for match in TOKEN_PATTERN.finditer(run.text):
                            full = match.group(0)
                            inner = match.group(1)
                            if full not in found:
                                found[full] = infer_token_type(inner)
    return [{"token": k, "type": v} for k, v in found.items()]


def populate_template(pptx_bytes: bytes, data_map: dict[str, str]) -> bytes:
    prs = Presentation(BytesIO(pptx_bytes))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for token, value in data_map.items():
                            if token in run.text:
                                run.text = run.text.replace(token, str(value))
    output = BytesIO()
    prs.save(output)
    return output.getvalue()
